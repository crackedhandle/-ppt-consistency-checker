import argparse
import json
import os
import re
import subprocess
import sys
import tempfile
from pathlib import Path
import logging

import google.generativeai as genai
from pdf2image import convert_from_path
from pptx import Presentation
import pytesseract

def setup_logging(debug=False):
    """Configure logging based on debug flag"""
    level = logging.DEBUG if debug else logging.INFO
    logging.basicConfig(
        format='%(asctime)s - %(levelname)s - %(message)s',
        level=level
    )

def extract_pptx_text(pptx_path, debug=False):
    """Extract text from PowerPoint slides and notes"""
    logging.info(f"Extracting text from PPTX: {pptx_path}")
    presentation = Presentation(pptx_path)
    slide_texts = []
    
    for i, slide in enumerate(presentation.slides):
        if debug:
            logging.debug(f"Processing slide {i+1}")
        slide_content = []
        
        # Extract slide title
        title = slide.shapes.title.text if slide.shapes.title else ""
        if title:
            slide_content.append(f"Title: {title}")
        
        # Extract content from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text.strip():
                        slide_content.append(text)
        
        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_content.append(f"Notes: {notes}")
        
        slide_text = "\n".join(slide_content)
        if debug and slide_text.strip():
            logging.debug(f"Slide {i+1} text:\n{slide_text}\n{'-'*40}")
        
        slide_texts.append(slide_text)
    
    return slide_texts

def convert_pptx_to_images(pptx_path, poppler_path=None, debug=False):
    """Convert PPTX to images using LibreOffice"""
    logging.info("Converting PPTX to images...")
    with tempfile.TemporaryDirectory() as tmp_dir:
        # Convert PPTX to PDF
        pdf_path = Path(tmp_dir) / "temp.pdf"
        try:
            # Convert all paths to strings for the command
            cmd = [
                "libreoffice", "--headless", "--convert-to", "pdf", 
                "--outdir", tmp_dir, str(pptx_path)
            ]
            if debug:
                # Convert all command elements to strings for logging
                cmd_str = " ".join(str(item) for item in cmd)
                logging.debug(f"Running command: {cmd_str}")
            
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True
            )
            
            if debug:
                if result.stdout:
                    logging.debug(f"LibreOffice stdout: {result.stdout}")
                if result.stderr:
                    logging.debug(f"LibreOffice stderr: {result.stderr}")
            
            result.check_returncode()
        except (subprocess.CalledProcessError, FileNotFoundError) as e:
            logging.warning(f"PPTX to PDF conversion failed: {str(e)}")
            if debug:
                if 'result' in locals():
                    logging.debug(f"Error details: {result.stderr}")
            return []
        
        # Convert PDF to images
        if pdf_path.exists():
            try:
                if debug:
                    logging.debug(f"Converting PDF to images with poppler_path={poppler_path}")
                
                images = convert_from_path(
                    str(pdf_path),
                    dpi=200,
                    poppler_path=poppler_path
                )
                return images
            except Exception as e:
                logging.error(f"PDF to image conversion error: {str(e)}")
                return []
        else:
            logging.warning(f"PDF not created: {pdf_path}")
            return []
    return []

def extract_image_text(image, debug=False):
    """Perform OCR on images using Tesseract"""
    try:
        text = pytesseract.image_to_string(image)
        if debug and text.strip():
            logging.debug(f"Extracted image text:\n{text}\n{'-'*40}")
        return text
    except Exception as e:
        logging.error(f"OCR Error: {str(e)}")
        return ""

def analyze_with_gemini(slide_texts, api_key, debug=False):
    """Detect inconsistencies using Gemini"""
    logging.info("Analyzing content with Gemini...")
    system_instruction = (
        "Analyze this presentation for inconsistencies. Return findings as JSON. "
        "Focus on these types:\n"
        "1. numerical: Conflicting numbers (revenue, percentages, statistics)\n"
        "2. textual: Contradictory claims or statements\n"
        "3. timeline: Mismatched dates, schedules, or forecasts\n"
        "4. logical: Reasoning flaws or contradictory conclusions\n\n"
        "Output format: List of objects with these keys:\n"
        "- slide_numbers: array of slide numbers involved (1-indexed)\n"
        "- description: clear explanation of the inconsistency\n"
        "- type: one of ['numerical', 'textual', 'timeline', 'logical']\n"
        "- confidence: float between 0-1\n\n"
        "Only include findings with confidence > 0.5. "
        "Reference slide numbers explicitly in descriptions."
    )
    
    # Build the slide content string
    parts = []
    for i, text in enumerate(slide_texts):
        parts.append(f"--- Slide {i+1} ---\n{text}")
    slide_content = "\n\n".join(parts)
    
    if debug:
        # Truncate for logging
        truncated = slide_content[:1000] + '... [truncated]' if len(slide_content) > 1000 else slide_content
        logging.debug(f"Sending to Gemini (truncated):\n{truncated}")
    
    try:
        model = genai.GenerativeModel(
            "gemini-1.5-flash-latest",
            system_instruction=system_instruction
        )
        
        response = model.generate_content(slide_content)
        
        if debug:
            logging.debug(f"Gemini raw response: {response.text}")
        
        return extract_json(response.text, debug)
    except Exception as e:
        logging.error(f"Gemini API error: {str(e)}")
        return {"error": str(e)}

def extract_json(text, debug=False):
    """Extract JSON from Gemini response with better handling of code blocks"""
    try:
        # Try to parse directly first
        return json.loads(text)
    except json.JSONDecodeError:
        try:
            # Handle code block format (```json ... ```)
            if text.strip().startswith('```'):
                # Extract content between the first and last triple backticks
                parts = text.split('```')
                if len(parts) >= 3:
                    json_content = parts[1]
                    # Remove optional "json" specifier
                    if json_content.lower().startswith('json'):
                        json_content = json_content[4:].lstrip()
                    return json.loads(json_content)
            
            # Try to find first JSON structure
            for i, char in enumerate(text):
                if char in {'[', '{'}:
                    return json.loads(text[i:])
            
            # If all else fails, try to clean the whole text
            cleaned = re.sub(r'[^\x00-\x7F]+', '', text)  # Remove non-ASCII
            return json.loads(cleaned)
        except Exception as e:
            if debug:
                logging.debug(f"JSON extraction failed: {str(e)}")
            return {"error": "JSON extraction failed", "raw_text": text}

def main():
    parser = argparse.ArgumentParser(
        description="PPTX Consistency Checker - Find inconsistencies in PowerPoint presentations"
    )
    parser.add_argument("pptx_path", help="Path to PowerPoint (.pptx) file")
    parser.add_argument("--output", default="results.json", help="Output JSON file path")
    parser.add_argument("--poppler", help="Path to Poppler bin directory")
    parser.add_argument("--api-key", help="Gemini API key")
    parser.add_argument("--debug", action="store_true", help="Enable debug logging")
    args = parser.parse_args()

    setup_logging(args.debug)
    logging.info("Starting PPTX consistency checker")

    # Validate inputs
    pptx_path = Path(args.pptx_path)
    if not pptx_path.exists():
        logging.error(f"File not found: {pptx_path}")
        sys.exit(1)
    
    # Set API key
    api_key = args.api_key or os.getenv("GOOGLE_API_KEY")
    if not api_key:
        logging.error("API key not provided. Use --api-key or set GOOGLE_API_KEY environment variable")
        sys.exit(1)
    genai.configure(api_key=api_key)
    
    # Configure Tesseract (use system default if not specified)
    if os.name == 'nt' and not os.getenv('TESSDATA_PREFIX'):
        logging.warning("TESSDATA_PREFIX environment variable not set - Tesseract might not work properly")

    # Process PPTX
    structured_text = extract_pptx_text(pptx_path, args.debug)
    
    # Convert to images and extract OCR text
    slide_images = convert_pptx_to_images(
        pptx_path, 
        poppler_path=args.poppler,
        debug=args.debug
    )
    
    ocr_texts = []
    if slide_images:
        logging.info(f"Extracting text from {len(slide_images)} slide images")
        for i, img in enumerate(slide_images):
            ocr_text = extract_image_text(img, args.debug)
            ocr_texts.append(ocr_text)
    else:
        logging.warning("No slide images generated - using only structured text")
        ocr_texts = [""] * len(structured_text)
    
    # Combine structured and OCR text
    combined_text = []
    for i, (structured, ocr) in enumerate(zip(structured_text, ocr_texts)):
        combined = f"{structured}\n\n[IMAGE CONTENT]\n{ocr}"
        combined_text.append(combined)
        if args.debug and combined.strip():
            logging.debug(f"Slide {i+1} combined text:\n{combined[:500]}...\n{'-'*40}")
    
    # Analyze with Gemini
    results = analyze_with_gemini(combined_text, api_key, args.debug)
    
    # Save results
    with open(args.output, "w") as f:
        json.dump(results, f, indent=2)
    
    logging.info(f"Analysis complete. Results saved to {args.output}")
    
    # Print summary
    if isinstance(results, list):
        logging.info(f"Found {len(results)} potential inconsistencies")
        for issue in results:
            logging.info(f"- Slide {issue['slide_numbers']}: {issue['description']} ({issue['type']}, confidence: {issue['confidence']:.2f})")
    else:
        logging.warning("No issues found or analysis failed. Check output file for details.")

if __name__ == "__main__":
    main()